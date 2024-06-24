import os
from pathlib import Path
from tqdm import tqdm
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Pt
from PIL import Image
import fitz


def create_dir(pdf_path:Path):
    '''
        Create a directory if it does not exist.
        args:
            pdf_path: Path, the directory path
        return:
            The directory path
    '''
    target_dir = pdf_path.parent/pdf_path.stem
    if not os.path.exists(target_dir):
        os.makedirs(target_dir)
    else:
        for f in os.listdir(target_dir):
            os.remove(f"{target_dir}{os.path.sep}{f}")
    return target_dir


def pdf2imglist(pdf_path:str, target_dir:str=None, dpi:int=300):
    '''
        Convert a PDF file to PNG files.
        args:
            pdf_path: str, the path to the PDF file
            target_dir: str, the directory to save the PNG files
            dpi: int, the resolution of the PNG files
    '''
    # Path to the PDF file
    pdf_path = Path(pdf_path)
    if target_dir is None:
        target_dir = create_dir(pdf_path)
        
    # Convert PDF to list of images
    images = convert_from_path(pdf_path, dpi=dpi)
    
    # Save images to files
    for i, image in tqdm(enumerate(images), total=len(images)):
        image.save(f'{target_dir}{os.path.sep}{pdf_path.stem}_{i}.png', 'PNG')
    print("Pdf was converted to images successfully.")



def imglist2pptx(image_folder_path, output_name='output.pptx'):
    '''
    Create a PowerPoint presentation from images in a folder.
        args:
            image_folder_path: str, Path to the folder containing images
            output_name: str, name of the output PowerPoint presentation
    '''
    image_folder_path = Path(image_folder_path)
    images  = os.listdir(image_folder_path)
    img = Image.open(image_folder_path/images[0])
    img_width, img_height = img.size

    slide_width = Pt(img_width*0.75)  # Pt = Px * 0.75
    slide_height = Pt(img_height*0.75)

    prs = Presentation()
    prs.slide_width  = slide_width
    prs.slide_height = slide_height
    left = top = Pt(0)

    for iimage in tqdm(images):
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.add_picture(
            str(image_folder_path/iimage), left, top, width=slide_width, height=slide_height)
    prs.save(image_folder_path.parent/output_name)
    print(f"PowerPoint presentation created as {output_name}")


def images2pdf(image_paths, output_pdf_path):
    # Create a new PDF
    doc = fitz.open()

    # Iterate over each image path
    for img_path in image_paths:
        # Add a new page with A4 size dimensions (595 x 842 points)
        page = doc.new_page(width = 595, height = 842)

        # Insert the image, fit the image to the page dimensions
        page.insert_image(page.rect, filename=str(img_path))

    # Save the resulting PDF
    doc.save(output_pdf_path)
    doc.close()

    print(f"PDF saved successfully to {output_pdf_path}")


def pdf2pptx(pdf_path, img_folder = None, output_name='output.pptx', dpi:int=300):
    '''
    Convert a PDF file to a PowerPoint presentation.
    args:
        pdf_path: str, Path to the PDF file
        output_name: str, name of the output PowerPoint presentation
    '''
    pdf_path = Path(pdf_path)
    if img_folder is None:
        img_folder = create_dir(pdf_path) 
    pdf2imglist(pdf_path, img_folder, dpi=dpi)
    imglist2pptx(img_folder, output_name)

# pdf_path = "./demo.pdf"
# pdf2pptx(pdf_path)

